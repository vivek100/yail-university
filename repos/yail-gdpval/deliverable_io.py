"""Read native professional deliverable formats for deterministic checks."""

from __future__ import annotations

import json
import re
from pathlib import Path
from typing import Any


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".xlsx":
        return _xlsx_text(path)
    if suffix == ".docx":
        return _docx_text(path)
    if suffix == ".pptx":
        return _pptx_text(path)
    if suffix == ".pdf":
        return _pdf_text(path)
    if suffix == ".ipynb":
        return _ipynb_text(path)
    if suffix in (".md", ".txt", ".csv", ".py", ".json"):
        try:
            return path.read_text(encoding="utf-8", errors="replace")
        except OSError:
            return ""
    return ""


def read_xlsx(path: Path) -> dict[str, list[list[Any]]]:
    wb = None
    try:
        from openpyxl import load_workbook

        wb = load_workbook(str(path), read_only=True, data_only=True)
    except Exception:
        return {}
    try:
        return {ws.title: [[cell for cell in row] for row in ws.iter_rows(values_only=True)] for ws in wb.worksheets}
    finally:
        if wb is not None:
            wb.close()


def read_pptx_slides(path: Path) -> list[str]:
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
    except Exception:
        return []

    slides: list[str] = []
    for slide in prs.slides:
        chunks: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    chunks.append("\t".join(cell.text for cell in row.cells))
        slides.append("\n".join(chunks))
    return slides


def approx_present(values: list[float], target: float, tol: float) -> bool:
    return any(abs(value - target) <= tol for value in values)


def numeric_values(sheets: dict[str, list[list[Any]]]) -> list[float]:
    nums: list[float] = []
    for rows in sheets.values():
        for row in rows:
            for cell in row:
                if isinstance(cell, bool):
                    continue
                if isinstance(cell, (int, float)):
                    nums.append(float(cell))
                elif isinstance(cell, str):
                    nums.extend(_numbers_from_text(cell))
    return nums


def _xlsx_text(path: Path) -> str:
    parts: list[str] = []
    for name, rows in read_xlsx(path).items():
        parts.append(f"## sheet: {name}")
        parts.extend("\t".join("" if cell is None else str(cell) for cell in row) for row in rows)
    return "\n".join(parts)


def _docx_text(path: Path) -> str:
    try:
        from docx import Document

        doc = Document(str(path))
    except Exception:
        return ""
    lines = [paragraph.text for paragraph in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            lines.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(lines)


def _pptx_text(path: Path) -> str:
    return "\n\n".join(f"## slide {idx + 1}\n{text}" for idx, text in enumerate(read_pptx_slides(path)))


def _pdf_text(path: Path) -> str:
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception:
        return ""


def _ipynb_text(path: Path) -> str:
    try:
        nb = json.loads(path.read_text(encoding="utf-8", errors="replace"))
    except (OSError, json.JSONDecodeError):
        return ""
    parts: list[str] = []
    for cell in nb.get("cells", []):
        src = cell.get("source", "")
        parts.append("".join(src) if isinstance(src, list) else str(src))
    return "\n".join(parts)


def _numbers_from_text(text: str) -> list[float]:
    out: list[float] = []
    for token in re.findall(r"-?\$?\(?\s*[\d,]+(?:\.\d+)?\s*\)?%?", text):
        cleaned = token.replace("$", "").replace(",", "").replace("%", "").strip()
        neg = cleaned.startswith("(") and cleaned.endswith(")")
        cleaned = cleaned.strip("()")
        if not cleaned or cleaned in {"-", "."}:
            continue
        try:
            value = float(cleaned)
        except ValueError:
            continue
        out.append(-value if neg else value)
    return out
